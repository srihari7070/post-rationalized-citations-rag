"""
Build the supervisor meeting slides.

Deliberately plain. Big type, few words, one idea per slide. The results carry
themselves and the questions need to be readable from across a table.

    python document/build_slides.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).resolve().parent / "Supervisor_Meeting_17Aug2026.pptx"

INK = RGBColor(0x1B, 0x21, 0x1F)
MUTE = RGBColor(0x6C, 0x75, 0x6F)
PINE = RGBColor(0x25, 0x5C, 0x4E)
RUST = RGBColor(0xA4, 0x50, 0x1C)
PAPER = RGBColor(0xF7, 0xF8, 0xF5)
LINE = RGBColor(0xD6, 0xDA, 0xD2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SERIF = "Georgia"
SANS = "Calibri"


def deck():
    p = Presentation()
    p.slide_width, p.slide_height = Inches(13.333), Inches(7.5)
    return p


def blank(prs, bg=PAPER):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bgf = s.background.fill
    bgf.solid()
    bgf.fore_color.rgb = bg
    return s


def tb(slide, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def para(tf, text, size, colour=INK, bold=False, font=SANS, space=10, first=False,
         align=PP_ALIGN.LEFT):
    par = tf.paragraphs[0] if first else tf.add_paragraph()
    par.alignment = align
    par.space_after = Pt(space)
    r = par.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = colour
    r.font.name = font
    return par


def rule(slide, x, y, w, colour=PINE, h=0.035):
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = colour
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def heading(slide, kicker, title, sub=None):
    rule(slide, 0.9, 0.72, 1.1)
    t = tb(slide, 0.9, 0.85, 11.5, 0.3)
    para(t, kicker.upper(), 12, PINE, True, SANS, 4, first=True)
    t2 = tb(slide, 0.9, 1.22, 11.5, 0.85)
    para(t2, title, 34, INK, False, SERIF, 6, first=True)
    if sub:
        t3 = tb(slide, 0.9, 2.15, 11.0, 0.6)
        para(t3, sub, 16, MUTE, False, SANS, 0, first=True)


def table(slide, rows, x, y, w, col_w=None, head=True, hi_row=None, size=14):
    n_r, n_c = len(rows), len(rows[0])
    h = 0.42 * n_r
    shp = slide.shapes.add_table(n_r, n_c, Inches(x), Inches(y), Inches(w), Inches(h))
    t = shp.table
    t.first_row = head
    if col_w:
        for i, cw in enumerate(col_w):
            t.columns[i].width = Emu(int(Inches(cw)))
    for ri, row in enumerate(rows):
        t.rows[ri].height = Inches(0.42)
        for ci, val in enumerate(row):
            c = t.cell(ri, ci)
            c.text = str(val)
            c.margin_left, c.margin_right = Inches(0.14), Inches(0.1)
            c.margin_top, c.margin_bottom = Inches(0.04), Inches(0.04)
            pr = c.text_frame.paragraphs[0]
            pr.alignment = PP_ALIGN.RIGHT if (ci and str(val)[:1].isdigit() or
                                              str(val).startswith(("-", "−", "+", "~"))) \
                else PP_ALIGN.LEFT
            if not pr.runs:           # empty cell still needs a run to style
                pr.add_run().text = ""
            f = pr.runs[0].font
            f.size = Pt(size)
            f.name = SANS
            if ri == 0 and head:
                f.bold, f.color.rgb = True, WHITE
                c.fill.solid(); c.fill.fore_color.rgb = PINE
            else:
                f.color.rgb = INK
                c.fill.solid()
                if hi_row is not None and ri == hi_row:
                    c.fill.fore_color.rgb = RGBColor(0xE2, 0xEC, 0xE7)
                    f.bold = True
                else:
                    c.fill.fore_color.rgb = WHITE
    return t


def callout(slide, x, y, w, h, label, text, colour=PINE, fill=RGBColor(0xE2, 0xEC, 0xE7)):
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = colour; box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.24)
    tf.margin_top = tf.margin_bottom = Inches(0.16)
    para(tf, label.upper(), 11, colour, True, SANS, 4, first=True)
    para(tf, text, 14, INK, False, SANS, 0)
    return box


def question(slide, num, q, why=None, y=2.6):
    """Returns the y the next element can safely start at."""
    lines_q = -(-len(q) // 62)
    lines_w = -(-len(why) // 96) if why else 0
    h = lines_q * 0.34 + (lines_w * 0.24 + 0.08 if why else 0)
    rule(slide, 0.9, y + 0.06, 0.028, PINE, max(0.32, h - 0.08))
    t = tb(slide, 1.15, y - 0.06, 11.2, h + 0.12)
    para(t, q, 21, INK, True, SANS, 4, first=True)
    if why:
        para(t, why, 14, MUTE, False, SANS, 0)
    return y + h + 0.28


def node(slide, x, y, w, h, title, body=None, colour=PINE,
         fill=WHITE, muted=False):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = colour
    box.line.width = Pt(1.5 if not muted else 0.75)
    box.shadow.inherit = False
    box.adjustments[0] = 0.08
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.15)
    tf.margin_top = tf.margin_bottom = Inches(0.1)
    para(tf, title, 14, MUTE if muted else INK, True, SANS, 3, first=True,
         align=PP_ALIGN.CENTER)
    if body:
        para(tf, body, 11.5, MUTE, False, SANS, 0, align=PP_ALIGN.CENTER)
    return box


def arrow(slide, x, y, w, h, shape=MSO_SHAPE.RIGHT_ARROW, colour=PINE):
    a = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = colour
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def caption(slide, x, y, w, text, colour=MUTE, size=11, align=PP_ALIGN.CENTER, bold=False):
    t = tb(slide, x, y, w, 0.4)
    para(t, text, size, colour, bold, SANS, 0, first=True, align=align)


def qa_column(slide, x, y, w, label, pairs, colour=PINE,
              fill=RGBColor(0xE2, 0xEC, 0xE7)):
    band = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.42))
    band.fill.solid(); band.fill.fore_color.rgb = fill
    band.line.fill.background(); band.shadow.inherit = False
    bt = band.text_frame
    bt.margin_left = Inches(0.2); bt.margin_top = Inches(0.05)
    para(bt, label.upper(), 11.5, colour, True, SANS, 0, first=True)
    cy = y + 0.58
    for q, a in pairs:
        t = tb(slide, x + 0.05, cy, w - 0.1, 0.34)
        para(t, q, 13.5, colour, True, SANS, 0, first=True)
        cy += 0.3
        t2 = tb(slide, x + 0.05, cy, w - 0.1, 0.5)
        para(t2, a, 13, INK, False, SANS, 0, first=True)
        cy += 0.3 + 0.24 * (len(a) // 52)
    return cy


def answer_card(slide, x, y, w, h, label, head, tail=None, colour=PINE,
                fill=WHITE, tail_colour=None):
    """A quoted answer. If tail is given it renders in tail_colour to show
    which words are shared with the comparison text."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = colour; box.line.width = Pt(1.25)
    box.shadow.inherit = False; box.adjustments[0] = 0.06
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.margin_top = tf.margin_bottom = Inches(0.13)
    para(tf, label.upper(), 10.5, colour, True, SANS, 4, first=True)
    par = tf.add_paragraph(); par.space_after = Pt(0)
    r1 = par.add_run(); r1.text = head
    r1.font.size = Pt(13); r1.font.name = SANS; r1.font.color.rgb = INK
    if tail:
        r2 = par.add_run(); r2.text = " " + tail
        r2.font.size = Pt(13); r2.font.name = SANS
        r2.font.color.rgb = tail_colour or RUST
        r2.font.bold = True
    return box


# ---------------------------------------------------------------- slides

prs = deck()

# 1 title
s = blank(prs, RGBColor(0x1B, 0x21, 0x1F))
rule(s, 1.1, 2.5, 1.4, RGBColor(0x7F, 0xBF, 0xA8))
t = tb(s, 1.1, 2.75, 11.0, 2.2)
para(t, "Correcting post-rationalised citations in RAG", 40,
     RGBColor(0xF7, 0xF8, 0xF5), False, SERIF, 10, first=True)
para(t, "Progress and open questions", 22, RGBColor(0x8A, 0x94, 0x8C), False, SANS, 0)
t2 = tb(s, 1.1, 5.7, 11.0, 0.9)
para(t2, "Srihari Ananthan", 15, RGBColor(0xB8, 0xC1, 0xBA), False, SANS, 3, first=True)
para(t2, "Supervisor meeting, 17 August 2026", 13, RGBColor(0x8A, 0x94, 0x8C), False, SANS, 0)

# 2  the problem
s = blank(prs)
heading(s, "The problem", "A citation that points at a real document it never used")
t = tb(s, 0.9, 2.6, 11.5, 2.2)
para(t, "A RAG system retrieves documents, writes an answer, and prints a citation next to "
        "each claim. Often the model wrote the answer from what it already knew, then attached "
        "whichever retrieved document looked closest.", 17, INK, False, SANS, 14, first=True)
para(t, "The document is real. It really was retrieved. The topic matches. The claim is often "
        "true as well.", 17, INK, False, SANS, 14)
para(t, "What is false is the relationship the citation implies between the two.",
     19, PINE, True, SANS, 0)
callout(s, 0.9, 5.3, 5.6, 1.7, "Why it is hard to catch",
        "A made-up citation pointing at a document that does not exist is easy to spot. This is "
        "not that. Nothing on the surface is wrong, because the evidence you would need is not "
        "in the text at all.")
callout(s, 7.0, 5.3, 5.4, 1.7, "Why it matters",
        "In a prototype a decorative citation is a curiosity. In a system informing a diagnosis "
        "or a contract it is part of the audit trail, and a false audit trail invites trust "
        "nobody earned.", RUST, RGBColor(0xF5, 0xE7, 0xDC))

# 3  how the system works, the diagram
s = blank(prs)
heading(s, "The plan", "How the system works",
        "One loop is closed. The other deliberately is not.")
caption(s, 0.9, 2.35, 4.4, "THE LOOP THAT RUNS", PINE, 11, PP_ALIGN.LEFT, True)
node(s, 0.9, 2.75, 2.55, 0.95, "1.  Model answers",
     "5 documents in, answer with citations out")
arrow(s, 3.6, 3.05, 0.42, 0.35)
node(s, 4.17, 2.75, 2.55, 0.95, "2.  Audit",
     "delete each cited passage, ask again, compare")
arrow(s, 6.87, 3.05, 0.42, 0.35)
node(s, 7.44, 2.75, 2.55, 0.95, "3.  Tell the model",
     "“citations [1] [3] were not needed, revise”")
arrow(s, 10.14, 3.05, 0.42, 0.35)
node(s, 10.71, 2.75, 1.75, 0.95, "4.  Re-audit", "measure again")
caption(s, 0.9, 3.85, 11.5,
        "The difference between the audit at step 2 and the re-audit at step 4 is the result.",
        INK, 13)
caption(s, 0.9, 4.5, 4.4, "THE BRANCH THAT DOES NOT", RUST, 11, PP_ALIGN.LEFT, True)
arrow(s, 5.2, 4.45, 0.35, 0.42, MSO_SHAPE.DOWN_ARROW, RGBColor(0xC2, 0xC8, 0xBC))
node(s, 3.3, 4.95, 4.3, 0.9, "Discriminator reads the same material",
     "gives its own verdict on each citation",
     RGBColor(0xC2, 0xC8, 0xBC), RGBColor(0xF7, 0xF8, 0xF5), muted=True)
node(s, 8.0, 4.95, 4.46, 0.9, "Recorded. Never acted on.",
     "answers RQ3: can a model spot this by reading?",
     RUST, RGBColor(0xF5, 0xE7, 0xDC))
arrow(s, 7.72, 5.25, 0.25, 0.32, MSO_SHAPE.RIGHT_ARROW, RGBColor(0xC2, 0xC8, 0xBC))
callout(s, 0.9, 6.15, 11.5, 0.95, "Why the discriminator is kept out of the loop",
        "If its opinion decided what to correct, I could not measure its accuracy, because it "
        "would be grading work it had caused. Keeping it outside is what makes the accuracy "
        "figure mean anything. Section 3.8.")

# 4  the same thing, in questions
s = blank(prs)
heading(s, "The plan", "The same thing, in questions")
qa_column(s, 0.9, 2.4, 5.6, "Branch one · the loop", [
    ("Is this the actual method?", "Yes. This is the thesis."),
    ("Does it find a better citation?", "No. It removes the unsupported one."),
    ("Proof?", "Citation correctness is unchanged: 70.0% → 70.0%."),
    ("So what does it achieve?",
     "The model stops claiming sources it never used. It repairs the audit trail, not the answer."),
])
qa_column(s, 7.0, 2.4, 5.4, "Branch two · the discriminator", [
    ("What is it?", "A hypothesis test, not a mechanism."),
    ("What hypothesis?", "That a cheap verbal check could replace the expensive causal test."),
    ("Does it work?", "No. It caught 32 of 375 fake citations. 8.5%."),
    ("Why is failing good news?", "It proves the causal test is necessary, not just preferred."),
], RUST, RGBColor(0xF5, 0xE7, 0xDC))
callout(s, 0.9, 6.35, 11.5, 0.8, "The distinction to hold onto",
        "The loop improves honesty, not accuracy. The discriminator never corrects anything, "
        "it only gets graded.")

# 5  the example: the question and its sources
s = blank(prs)
heading(s, "One question, end to end", "Question A04 and what retrieval returned")
t = tb(s, 0.9, 2.35, 11.5, 0.75)
para(t, "“Which Swiss climate startup uses drones and environmental DNA tracking to replant "
        "mangroves at scale along coastlines threatened by erosion?”",
     18, INK, True, SERIF, 0, first=True)
caption(s, 0.9, 3.2, 11.5, "The 5 closest company profiles out of 38,692",
        MUTE, 12, PP_ALIGN.LEFT)
table(s, [
    ["", "Company", "What it does"],
    ["[1]", "Inverto Earth AG", "Bern. Drones and eDNA, mangrove restoration"],
    ["[2]", "DNAir AG", "Zürich area. eDNA sensors, but from air"],
    ["[3]", "Rrreefs", "Zürich. 3D-printed coral reefs"],
    ["[4]", "Open Climate Solutions", "Baar. Macroalgae for agriculture"],
    ["[5]", "Open Forest Protocol", "Lausanne. Blockchain for forest monitoring"],
], 0.9, 3.6, 11.5, col_w=[0.9, 3.4, 7.2], hi_row=1, size=13)
callout(s, 0.9, 6.5, 11.5, 0.7, "Note",
        "All five are Swiss climate startups. Only [1] does mangroves with drones. "
        "The recorded correct answer is Inverto Earth AG.")

# 6  the model answers
s = blank(prs)
heading(s, "One question, end to end", "Step 1  ·  Mistral answers")
answer_card(s, 0.9, 2.6, 11.5, 1.25, "The answer, with all five profiles in the prompt",
            "“Inverto Earth AG [1] is the Swiss climate startup that uses drones and "
            "environmental DNA tracking to replant mangroves at scale along coastlines "
            "threatened by erosion.”")
t = tb(s, 0.9, 4.15, 11.5, 1.0)
para(t, "Correct company. Cites [1]. Looks perfect.", 20, PINE, True, SANS, 10, first=True)
para(t, "Every surface check passes. The document is real, it was retrieved, the topic matches, "
        "and the company named is the right one.", 16, INK, False, SANS, 0)
callout(s, 0.9, 5.5, 11.5, 1.3, "But look at what is missing",
        "The answer says nothing that only profile [1] knows. No Bern. No Hilterfingen. No "
        "detail about how the planting works. It is the question read back with a name attached "
        "to the front.", RUST, RGBColor(0xF5, 0xE7, 0xDC))

# 7  the test
s = blank(prs)
heading(s, "One question, end to end", "Step 2  ·  Delete [1] and ask again")
HEAD_A = "Inverto Earth AG [1] is the"
HEAD_C = "The provided sources do not mention any"
TAIL   = ("Swiss climate startup that uses drones and environmental DNA tracking to replant "
          "mangroves at scale along coastlines threatened by erosion.")
answer_card(s, 0.9, 2.5, 11.5, 1.3, "The original answer", HEAD_A, TAIL)
answer_card(s, 0.9, 4.0, 11.5, 1.3, "The answer with profile [1] deleted", HEAD_C, TAIL,
            RGBColor(0xC2, 0xC8, 0xBC), RGBColor(0xF7, 0xF8, 0xF5))
t = tb(s, 0.9, 5.5, 5.6, 0.9)
para(t, "The orange text is word for word identical.", 15, RUST, True, SANS, 5, first=True)
para(t, "20 of the 26 words. 77% of the answer.", 15, INK, False, SANS, 0)
callout(s, 7.0, 5.45, 5.4, 1.35, "Verdict",
        "Similarity 0.8541, above the 0.85 cutoff. The answer survived without the source, so "
        "the citation was not doing the work. Post-rationalised.",
        RUST, RGBColor(0xF5, 0xE7, 0xDC))

# 8  tell it, it rewrites
s = blank(prs)
heading(s, "One question, end to end", "Steps 3 and 4  ·  Tell it, and it rewrites")
callout(s, 0.9, 2.5, 11.5, 0.95, "What we send back",
        "“You previously answered the question below and cited source [1]. Testing showed "
        "that removing it did not change your answer. It was not genuinely needed. Revise.”")
answer_card(s, 0.9, 3.75, 11.5, 1.6, "The revised answer",
            "“Inverto Earth AG, a Swiss climate tech company based in ",
            "Bern, Hilterfingen, leverages drone-based planting systems and environmental DNA "
            "tracking to efficiently restore and monitor mangrove ecosystems at scale along "
            "coastlines threatened by erosion [1].”", PINE, WHITE, PINE)
callout(s, 0.9, 5.55, 11.5, 1.3, "Everything in green came out of profile [1]",
        "Bern. Hilterfingen. Drone-based planting systems. Restore and monitor mangrove "
        "ecosystems. None of those words appears in the question, and none appears in any of "
        "the other four profiles.")

# 9  test again
s = blank(prs)
heading(s, "One question, end to end", "Step 5  ·  Delete [1] again")
t = tb(s, 0.9, 2.4, 11.5, 0.5)
para(t, "Same deletion. The model produces the byte-identical refusal it gave in step 2.",
     16, MUTE, False, SANS, 0, first=True)
table(s, [
    ["", "Compared against", "Similarity", "Verdict"],
    ["Before", "the question restated, with a name attached", "0.8541", "post-rationalised"],
    ["After", "facts that only profile [1] contains", "0.8183", "genuine"],
], 0.9, 3.0, 11.5, col_w=[1.5, 5.6, 2.2, 2.2], hi_row=2, size=15)
t2 = tb(s, 0.9, 4.65, 11.5, 1.2)
para(t2, "The deleted passage was the same. The refusal was the same. What changed is the "
         "answer being compared to it.", 17, INK, False, SANS, 10, first=True)
para(t2, "Before, the answer sat almost on top of the refusal, because both were just the "
         "question. After, it moved away, because it now carried facts that vanish when you "
         "delete [1].", 16, MUTE, False, SANS, 0)
callout(s, 0.9, 6.2, 11.5, 0.85, "What the loop actually did",
        "It did not change which company was cited. It forced the answer to be built out of the "
        "source it was already pointing at.")

# 10  the 75 questions
s = blank(prs)
heading(s, "At scale", "The 75 test questions",
        "That was one. Each type exists to catch a different behaviour.")
table(s, [
    ["Type", "n", "What it is", "Why it exists"],
    ["A  Answerable", "30", "Exactly one company fits", "Cleanest test: does it cite what it used?"],
    ["B  Ambiguous", "15", "Many companies half-fit", "Does it hedge by citing everything?"],
    ["C  Hard", "20", "Answer exists, wording avoids the profile", "Does retrieval difficulty change citing?"],
    ["D  Unanswerable", "10", "Nothing in the corpus fits", "What does it do with no right answer?"],
], 0.9, 2.9, 11.5, col_w=[2.3, 0.7, 4.0, 4.5])
callout(s, 0.9, 5.6, 11.5, 1.45, "A failure worth mentioning",
        "My first question set was generic, e.g. “which companies work on "
        "sustainability”. Post-rationalisation came out above 90% everywhere. The "
        "measurement was fine, the questions were broken: when any of the five documents could "
        "support the answer, removing one changes nothing. A04 works as a test precisely "
        "because only one profile can answer it.",
        RUST, RGBColor(0xF5, 0xE7, 0xDC))

# 11  result 1
s = blank(prs)
heading(s, "Result 1", "Correction works, but only on some models")
t0 = tb(s, 0.9, 2.45, 11.5, 0.45)
para(t0, "Detection works on all three. Repair does not.", 17, PINE, True, SANS, 0, first=True)
table(s, [
    ["Model", "Bad citations found", "PRR before", "PRR after", "Change", "Repairs when told"],
    ["Mistral 7B", "33 of 106", "22.2%", "9.4%", "−12.8pp", "59.1%"],
    ["Llama 3 8B", "35 of 120", "20.6%", "14.1%", "−6.4pp", "36.4%"],
    ["Gemini 2.5 Flash", "57 of 131", "25.0%", "23.8%", "−1.2pp  (n.s.)", "14.8%"],
], 0.9, 3.0, 11.5, col_w=[2.5, 2.1, 1.7, 1.6, 1.9, 1.7], hi_row=1, size=14)
t = tb(s, 0.9, 4.9, 11.5, 0.7)
para(t, "The audit finds plenty in Gemini’s answers, more than in either local model. The "
        "step that fails is the repair: told which citations failed, Gemini does not remove "
        "them.", 15, INK, False, SANS, 0, first=True)
callout(s, 0.9, 5.7, 11.5, 1.3, "Rewriting is not fixing",
        "Gemini rewrote 25 of its 75 answers when asked, so it read the instruction and produced "
        "different text. It kept the bad citations anyway. That is why I measure whether the "
        "problem got fixed, not whether the answer changed.")

# 12  result 2, the strongest
s = blank(prs)
heading(s, "Result 2", "The discriminator cannot spot a fake citation by reading")
t = tb(s, 0.9, 2.5, 11.5, 0.45)
para(t, "Accuracy looks fine at 59.5% to 72.9%. It is not.", 19, INK, True, SANS, 0, first=True)
table(s, [
    ["Discriminator", "Fakes caught", "Missed", "Recall", "Accuracy"],
    ["Gemini 2.5 Flash", "10", "117", "7.9%", "63.3%"],
    ["Mistral 7B", "0", "123", "0.0%", "65.5%"],
    ["Llama 3 8B", "22", "103", "17.6%", "66.7%"],
    ["All nine pairings", "32", "343", "8.5%", ""],
], 0.9, 3.1, 11.5, col_w=[3.0, 2.3, 2.0, 2.1, 2.1], hi_row=2, size=14)
t2 = tb(s, 0.9, 5.3, 11.5, 0.55)
para(t2, "A discriminator answering “genuine” to everything would score 65.0%, right "
         "in the middle of that range, while catching none of the 375.",
     16, INK, True, SANS, 0, first=True)
callout(s, 0.9, 6.15, 11.5, 0.95, "Why this is the strongest thing I have",
        "It does not depend on my threshold, on how I average, or on the validation study that "
        "failed. It survives every objection to the rest of the thesis, and it is what justifies "
        "the expensive causal test over a cheap verbal one.")

# 13  result 3
s = blank(prs)
heading(s, "Result 3", "Which model does the judging makes no difference")
table(s, [
    ["Generator \\ Discriminator", "Gemini", "Mistral", "Llama 3"],
    ["Gemini", "−3.5", "+0.0", "−0.2"],
    ["Mistral", "−12.8", "−12.8", "−12.8"],
    ["Llama 3", "−6.4", "−6.4", "−6.4"],
], 2.6, 2.75, 8.1, col_w=[3.0, 1.7, 1.7, 1.7], size=16)
t = tb(s, 0.9, 5.0, 11.5, 1.0)
para(t, "The discriminators do disagree with each other, on 9 and 13 of 75 questions. The "
        "revised answers come out identical on all 75 anyway. This is by design: the "
        "discriminator never steers the loop.", 15, INK, False, SANS, 0, first=True)
callout(s, 0.9, 6.05, 11.5, 1.0, "Turned into an advantage",
        "If the discriminator cannot affect the outcome, three conditions per model are the same "
        "experiment run three times. Their spread measures my noise for free. That is how I know "
        "Gemini’s apparent −3.5 is not real.")

# 14  honesty vs accuracy, part one
s = blank(prs)
heading(s, "What this means", "Honesty and accuracy are two different things",
        "The thesis measures both, separately, on purpose.")
table(s, [
    ["", "PRR  ·  honesty", "CCR  ·  accuracy"],
    ["The question it asks", "Did you actually read the document you pointed at?",
     "Did you point at the right document?"],
    ["How it is measured", "Delete the passage, regenerate, compare",
     "String match against the recorded right answer"],
    ["When", "During the run, live", "Afterwards, from the logs"],
    ["Cost", "Extra model calls", "Free"],
], 0.9, 2.9, 11.5, col_w=[2.6, 4.6, 4.3], size=13)
t = tb(s, 0.9, 5.35, 11.5, 0.6)
para(t, "CCR is checked on the 50 questions that have a recorded right answer, types A, B and "
        "C. Type D has no right answer, so it is scored separately as abstention.",
     14, MUTE, False, SANS, 0, first=True)
callout(s, 0.9, 6.05, 11.5, 1.0, "You can be right and dishonest",
        "A model can name the correct company without having used the profile, which scores well "
        "on CCR and badly on PRR. It can also faithfully use a profile that names the wrong "
        "company, which scores the other way round.")

# 15  honesty vs accuracy, part two
s = blank(prs)
heading(s, "What this means", "The loop improves honesty. It does not improve accuracy.")
table(s, [
    ["Model", "CCR before", "CCR after", "PRR before", "PRR after"],
    ["Mistral 7B", "70.0%", "70.0%", "22.2%", "9.4%"],
    ["Gemini 2.5 Flash", "76.0%", "76.0%", "25.0%", "23.8%"],
    ["Llama 3 8B", "74.0%", "72.0%", "20.6%", "14.1%"],
], 0.9, 2.6, 11.5, col_w=[2.9, 2.15, 2.15, 2.15, 2.15], hi_row=1, size=14)
caption(s, 0.9, 4.25, 11.5, "Accuracy does not move. Honesty does.", INK, 17, PP_ALIGN.LEFT, True)
t = tb(s, 0.9, 4.75, 5.6, 1.6)
para(t, "Back to A04", 15, PINE, True, SANS, 5, first=True)
para(t, "Before:  cited Inverto Earth AG, correct, PRR 100%", 14, INK, False, SANS, 3)
para(t, "After:   cited Inverto Earth AG, correct, PRR 0%", 14, INK, False, SANS, 6)
para(t, "Same company both times. Correctness never moved.", 14, MUTE, False, SANS, 0)
callout(s, 7.0, 4.7, 5.4, 2.15, "The answer to give if asked",
        "“Does your system make the answers better?” No. Same companies cited, same "
        "70% correct. What it does is make the citation mean something. Before, the bracket was "
        "decoration you could not trust. After, following it back gets you to the text that "
        "actually produced the sentence. The loop does not improve the answer. It repairs the "
        "audit trail.")

# 16  honest limits
s = blank(prs)
heading(s, "What is not safe", "Two things I raise before anyone asks")
callout(s, 0.9, 2.6, 5.6, 2.0, "1. The effect peaks at my threshold",
        "Mistral’s −12.8pp becomes −1.1 at a cutoff of 0.90 and zero at 0.95. I "
        "chose 0.85 by convention, not from data. The direction holds everywhere. The size does "
        "not.", RUST, RGBColor(0xF5, 0xE7, 0xDC))
callout(s, 6.9, 2.6, 5.5, 2.0, "2. The measure is unvalidated",
        "I built a blind study to check it against human judgment and ran it twice. Both failed. "
        "No published work using this method has validated it either.",
        RUST, RGBColor(0xF5, 0xE7, 0xDC))
t = tb(s, 0.9, 4.9, 11.5, 2.0)
para(t, "What the failed validation did establish", 17, PINE, True, SANS, 6, first=True)
para(t, "“The model used this passage” has three defensible meanings, giving 0%, 10.2% "
        "and 20 to 25% on identical data. They differ in granularity. I measure the coarsest and "
        "had not said so.", 15, INK, False, SANS, 8)
para(t, "One consequence is provable without any annotation: whole-answer similarity cannot see "
        "the loss of one item from a list. Across 1,031 cases it misjudges 12% of the time when "
        "the answer names one company, rising to 82% when it names five.",
     15, INK, False, SANS, 0)

# 8 questions: length
s = blank(prs)
heading(s, "Questions", "How long should the thesis be?")
y = question(s, 1, "What is the word count, and is there a limit?",
         "The template gives no number. It only says the abstract is 300 words maximum.", 2.7)
y = question(s, 2, "Or is it counted in pages rather than words?", None, y)
y = question(s, 3, "Does the count include tables, figures and the reference list?", None, y)
t = tb(s, 0.9, y + 0.25, 11.5, 1.4)
para(t, "Where I am now", 15, PINE, True, SANS, 5, first=True)
para(t, "Six chapters, about 17,200 words, 19 tables, three appendices. If the expectation is "
        "22,000 to 25,000 I have room to expand. If it is 80 pages I need to know whether that "
        "is with or without appendices.", 15, INK, False, SANS, 0)

# 9 questions: appendix and code
s = blank(prs)
heading(s, "Questions", "Appendix, code, images and results")
y = question(s, 1, "How do I share the code?",
         "A link to a repository, a zip with the submission, or printed in an appendix?", 2.6)
y = question(s, 2, "Where do the results files go?",
         "I have 12 run logs, roughly 900 answers with their audits. Too large to print.", y)
y = question(s, 3, "How should figures and tables be numbered and captioned?",
         "The template automates captions. Should I use its numbering, or is mine fine?", y)
y = question(s, 4, "Does the appendix count towards the length?", None, y)
callout(s, 0.9, y + 0.15, 11.5, 0.85, "What I have ready",
        "Appendix A: the twelve conditions.  Appendix B: the annotation codebook.  "
        "Appendix C: where every script and log lives.")

# 10 questions: references
s = blank(prs)
heading(s, "Questions", "Are 22 references enough?")
t = tb(s, 0.9, 2.7, 11.5, 1.15)
para(t, "I checked this before asking. The literature review cites 22 distinct works, every "
        "one verified against the published record, and all 22 are cited in the text rather "
        "than padding a list.", 16, INK, False, SANS, 10, first=True)
y = question(s, 1, "Is 22 enough for an MSc literature review?", None, 4.2)
y = question(s, 2, "What final number do you expect?",
         "I have around 160 papers screened as worth reading in full, so I can go further.", y)
y = question(s, 3, "Should the systematic search be reported formally?",
         "I have PRISMA numbers: 2,224 records found, 1,811 screened, 596 kept. Chapter 2 "
         "currently spends about 500 words on the search. Cut it, or keep it?", y)

# 11 questions: the big two
s = blank(prs)
heading(s, "Questions", "The two that change what I do next")
callout(s, 0.9, 2.7, 11.5, 1.85, "1. What should “the model used this source” actually mean?",
        "My number is 4% or 25% depending on which definition we pick, and it is the thesis’s "
        "main measurement. I need one definition, stated in words before any formula, and a "
        "reason for choosing it. I do not think I should decide this alone.",
        RUST, RGBColor(0xF5, 0xE7, 0xDC))
callout(s, 0.9, 4.8, 11.5, 1.7, "2. My significance test is the wrong test",
        "I used an unpaired test on paired data, across nine conditions, with no correction for "
        "multiple comparisons. The correct paired test gives p = 0.0006 instead of 0.031, which "
        "survives correction. My main result is stronger than I claimed, and I want to confirm "
        "the fix before rewriting.", RUST, RGBColor(0xF5, 0xE7, 0xDC))
t = tb(s, 0.9, 6.65, 11.5, 0.6)
para(t, "Both are written up in full in QUESTIONS_FOR_PROFESSOR.md if you would like the detail.",
     14, MUTE, False, SANS, 0, first=True)

# 12 close
s = blank(prs, RGBColor(0x1B, 0x21, 0x1F))
rule(s, 1.1, 2.6, 1.4, RGBColor(0x7F, 0xBF, 0xA8))
t = tb(s, 1.1, 2.85, 11.0, 3.2)
para(t, "Where I stand", 32, RGBColor(0xF7, 0xF8, 0xF5), False, SERIF, 16, first=True)
for line in [
    "Correction works on two of three models, and I can measure why.",
    "Models cannot detect this by reading. That is the robust finding.",
    "The adversarial framing does not fit what I built. I should drop the word.",
    "The measure needs a definition before it can be validated.",
]:
    para(t, "—   " + line, 17, RGBColor(0xE8, 0xEC, 0xE7), False, SANS, 11)

prs.save(str(OUT))
print(f"{OUT.name}  —  {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
